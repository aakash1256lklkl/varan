"""
Varan PowerPoint editor — create, edit, read and summarize .pptx files
built on python-pptx. Supports slides, textboxes, tables and charts.
"""
from __future__ import annotations

from pathlib import Path
from typing import Any, Optional

from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from . import ppt_live as _live


# Literal filler tokens agents sometimes emit instead of real content. Writing
# these into a deck is a silent failure (the file "saves ok" but contains
# placeholder garbage like "- item"); we refuse them so the model must supply
# REAL content or use edit_slides to fill an existing text box.
_FILLER_TOKENS = {
    "item", "items", "item1", "item 1", "insert", "placeholder", "placeholders",
    "your text here", "text here", "text", "content", "contents", "sample",
    "sample text", "lorem", "ipsum", "lorem ipsum", "todo", "tbd", "t.b.d.",
    "n/a", "na", "- item", "-item", "• item", "bullet", "bullets",
}


def _is_filler(value) -> bool:
    """True when a cell/bullet is a degenerate placeholder rather than content."""
    t = str(value or "").strip().lower()
    return t in _FILLER_TOKENS


def _assert_real_slide_content(spec: dict):
    """Refuse slides whose body/table consists purely of placeholder filler.

    Guards create_presentation / rebuild_slides / add_slides / edit_slides
    against silently writing '- item' bullets or 'item' / 'i | t | e | m'
    tables. Errored slides are never written: load/save happens as a whole
    transaction, so a refusal leaves the deck untouched.
    """
    bullets = spec.get("bullets")
    if bullets is not None:
        real = [b for b in bullets if (b or "").strip()]
        if real and all(_is_filler(b) for b in real):
            raise ValueError(
                "Refusing placeholder-only content: every bullet is filler "
                f"({real!r}). Provide REAL slide content, or use "
                "edit_slides with 'set_bullets' / 'add_table' to fill an "
                "existing text box."
            )
    table = spec.get("table")
    if table:
        headers = [h for h in (table.get("headers") or []) if (h or "").strip()]
        rows = [[v for v in (r or []) if (v or "").strip()] for r in (table.get("rows") or [])]
        cells = [v for row in rows for v in row]
        if headers and all(_is_filler(h) for h in headers) and (
            not cells or all(len(str(v).strip()) <= 1 for v in cells)
        ):
            raise ValueError(
                "Refusing placeholder-only table: "
                f"headers={headers!r} rows={[[str(v) for v in r] for r in rows]!r}. "
                "Provide REAL table data (headers + cell values), or use "
                "edit_slides with 'add_table' to fill a table on an existing slide."
            )


class PowerPointEditor:
    extension = ".pptx"

    # ------------------------------------------------------------------
    # Creation
    # ------------------------------------------------------------------
    def create_presentation(self, path: str | Path,
                            slides: Optional[list[dict]] = None) -> str:
        """Create a new presentation.

        slide dicts:
          {"layout": "title"|"bullets"|"blank", "title": "...", "subtitle": "...",
           "bullets": [...], "table": {"headers": [...], "rows": [[...]]},
           "chart": {"type": "bar", "title": "...", "categories": [...], "data": [...]}}
        """
        path = Path(path)
        path.parent.mkdir(parents=True, exist_ok=True)
        prs = Presentation()

        for slide in (slides or []):
            self._add_slide(prs, slide)

        prs.save(str(path))
        return str(path)

    def _add_slide(self, prs: Presentation, spec: dict):
        _assert_real_slide_content(spec)
        layout_name = (spec.get("layout") or "bullets").lower()
        layout = self._layout(prs, layout_name)
        slide = prs.slides.add_slide(layout)

        title = spec.get("title")
        if title is not None and slide.shapes.title is not None:
            slide.shapes.title.text = title

        subtitle = spec.get("subtitle")
        if subtitle:
            if layout_name == "title":
                # put subtitle in the subtitle placeholder if present
                for ph in slide.placeholders:
                    if ph.placeholder_format.idx == 1:
                        ph.text = subtitle
                        break
                else:
                    self._add_textbox(slide, subtitle, top=3.0, height=1.0)
            else:
                self._add_textbox(slide, subtitle, top=1.2, height=0.8)

        bullets = spec.get("bullets")
        if bullets:
            body_idx = 1
            for ph in slide.placeholders:
                if ph.placeholder_format.type == 2:  # BODY
                    tf = ph.text_frame
                    first = True
                    for b in bullets:
                        p = tf.paragraphs[0] if first else tf.add_paragraph()
                        first = False
                        p.text = str(b)
                        p.level = 0
                    break
            else:
                # No body placeholder -> add a textbox
                self._add_textbox(slide, "", top=1.5, height=4.5, bullets=bullets)

        table = spec.get("table")
        if table:
            self._add_table(slide, table)

        chart = spec.get("chart")
        if chart:
            self._add_chart(slide, chart)

    def _edit_one_slide(self, slide, spec: dict):
        """Apply a surgical text/shape edit to one existing slide.

        Supported keys:
          title:        replace the slide's title text
          set_bullets:  replace the body placeholder (or first textbox) bullets
          replace:      dict {old_text: new_text} applied to every text shape
          add_textbox:  {"text": "...", "left"/"top"/"width"/"height"} — add a
                        new text box to THIS slide (not a new slide)
          add_table:    {"headers": [...], "rows": [[...]]} — add a table to
                        THIS slide
          add_chart:    existing chart spec — add a chart to THIS slide
        """
        title = spec.get("title")
        if title is not None and slide.shapes.title is not None:
            slide.shapes.title.text = title

        _assert_real_slide_content({
            "bullets": spec.get("set_bullets"),
            "table": spec.get("add_table"),
        })

        replace_map = spec.get("replace") or {}
        if replace_map:
            all_shapes = list(slide.shapes)
            table_shapes = []
            for shape in all_shapes:
                if shape.has_table:
                    table_shapes.append(shape)
            # text runs
            for shape in all_shapes:
                if not shape.has_text_frame:
                    continue
                tf = shape.text_frame
                for para in tf.paragraphs:
                    for run in para.runs:
                        new_text = run.text
                        for old_text, new_sub in replace_map.items():
                            if old_text and old_text in new_text:
                                new_text = new_text.replace(old_text, new_sub)
                        run.text = new_text
            # table cells (complex edit: replace inside tables too)
            for shape in table_shapes:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if not cell.text_frame:
                            continue
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                new_text = run.text
                                for old_text, new_sub in replace_map.items():
                                    if old_text and old_text in new_text:
                                        new_text = new_text.replace(old_text, new_sub)
                                run.text = new_text

        bullets = spec.get("set_bullets")
        if bullets is not None:
            body_ph = None
            for ph in slide.placeholders:
                if ph.placeholder_format.type == 2:  # BODY
                    body_ph = ph
                    break
            tf = body_ph.text_frame if body_ph is not None else None
            if tf is None:
                # fall back to the first text box that is not the title
                for shape in slide.shapes:
                    if shape.has_text_frame and shape is not slide.shapes.title:
                        tf = shape.text_frame
                        break
            if tf is not None:
                tf.clear()
                first = True
                for b in bullets:
                    p = tf.paragraphs[0] if first else tf.add_paragraph()
                    first = False
                    p.text = str(b)
                    p.level = 0

        add_tb = spec.get("add_textbox")
        if add_tb:
            d = add_tb if isinstance(add_tb, dict) else {}
            self._add_textbox(
                slide,
                str(d.get("text", "") if isinstance(add_tb, dict) else (add_tb or "")),
                left=float(d.get("left", 1.0)),
                top=float(d.get("top", 1.5)),
                width=float(d.get("width", 8.5)),
                height=float(d.get("height", 1.0)),
            )

        add_table = spec.get("add_table")
        if add_table is not None:
            self._add_table(slide, add_table)

        add_chart = spec.get("add_chart")
        if add_chart is not None:
            self._add_chart(slide, add_chart)

    def _remove_slide(self, prs, index: int):
        """Delete one slide (1-based) from the deck in place, dropping its slide
        relationship so the saved file no longer contains it. This is the
        surgical 'remove slide N' operation (no full rebuild needed)."""
        sldIdLst = prs.slides._sldIdLst  # noqa: SLF001
        if index < 1 or index > len(sldIdLst):
            raise IndexError(f"slide index out of range: {index}")
        sldId = sldIdLst[index - 1]
        rId = sldId.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        if rId:
            prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)

    def _layout(self, prs, name):
        for layout in prs.slide_layouts:
            if layout.name and layout.name.lower() == name:
                return layout
        # fallbacks
        mapping = {
            "title": 0, "bullets": 1, "blank": 6,
        }
        idx = mapping.get(name, 1)
        return prs.slide_layouts[idx]

    def _add_textbox(self, slide, text, left=1.0, top=1.0, width=8.5, height=4.0, bullets=None):
        from pptx.enum.text import PP_ALIGN
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        if text:
            tf.text = text
        if bullets:
            first = True
            for b in bullets:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.text = str(b)
                p.level = 0
        return box

    def _add_table(self, slide, spec):
        headers = spec.get("headers") or []
        rows = spec.get("rows") or []
        rows_count = len(rows) + (1 if headers else 0)
        cols_count = max(len(headers), *(len(r) for r in rows)) if rows else len(headers)
        if rows_count == 0 or cols_count == 0:
            return
        shape = slide.shapes.add_table(
            rows_count, cols_count, Inches(0.6), Inches(2.0), Inches(8.8), Inches(3.0)
        )
        table = shape.table
        ri = 0
        if headers:
            for ci, h in enumerate(headers):
                table.cell(0, ci).text = str(h)
            ri = 1
        for r in rows:
            vals = list(r)
            for ci in range(cols_count):
                v = vals[ci] if ci < len(vals) else ""
                table.cell(ri, ci).text = str(v)
            ri += 1

    def _add_chart(self, slide, spec):
        try:
            from pptx.chart.data import CategoryChartData
            from pptx.enum.chart import XL_CHART_TYPE
        except ImportError:
            return
        ctype = (spec.get("type") or "bar").lower()
        categories = spec.get("categories") or []
        data = spec.get("data") or []
        if not categories or not data:
            return
        chart_data = CategoryChartData()
        chart_data.categories = categories
        # data could be a single list (one series) -> assume "values"
        series_name = spec.get("series_name", "Values")
        chart_data.add_series(series_name, data)
        if ctype == "line":
            xl = XL_CHART_TYPE.LINE_MARKERS
        elif ctype == "pie":
            xl = XL_CHART_TYPE.PIE
        else:
            xl = XL_CHART_TYPE.COLUMN_CLUSTERED
        graphic_frame = slide.shapes.add_chart(
            xl, Inches(1.0), Inches(2.5), Inches(6), Inches(4), chart_data
        )
        if spec.get("title"):
            graphic_frame.chart.has_title = True
            graphic_frame.chart.chart_title.text_frame.text = spec["title"]

    # ------------------------------------------------------------------
    # Editing
    # ------------------------------------------------------------------
    def edit_presentation(self, path: str | Path,
                          add_slides: Optional[list[dict]] = None,
                          rebuild_slides: Optional[list[dict]] = None,
                          edit_slides: Optional[list[dict]] = None,
                          remove_slides: Optional[list[int]] = None,
                          inplace: bool = False) -> "tuple[str, dict]":
        """Edit an existing presentation.

        Modes (combinable except rebuild, which is exclusive):
          * append: `add_slides` are appended to the existing deck.
          * rebuild: `rebuild_slides` REPLACE every existing slide with a fresh
            deck built from that list (used to retheme / rewrite a deck around
            a new idea instead of just stacking slides on).
          * surgical: `edit_slides` changes the text/shape content of SPECIFIC
            existing slides without touching the rest of the deck. Each item is:
              {"slide": 1-based_index,
               "title": "new title",            # replace this slide's title
               "set_bullets": [...],            # replace the body bullets
               "replace": {"old": "new", ...},  # find/replace text (runs + table cells)
               "add_textbox": {"text": ...},    # add a textbox to this slide
               "add_table": {"headers": [...], "rows": [[...]]},
               "add_chart": {chart-spec}}
          * remove: `remove_slides` deletes SPECIFIC slides (1-based indices)
            from the deck — the surgical 'delete slide N' without a rebuild.
            Removals are applied after edits, against the ORIGINAL deck
            numbering.

        By default writes a "NAME_edited.pptx" copy; inplace=True saves back
        over the original (used when editing the selected target file).

        Returns (path_written, summary) where summary reports exactly how many
        slides were present before, how many were appended/rebuilt/edited/
        removed, so callers can describe the true outcome.
        """
        src = Path(path)
        if not src.exists():
            raise FileNotFoundError(src)
        dst = src if inplace else src.with_name(f"{src.stem}_edited{self.extension}")

        before = 0
        try:
            prs = Presentation(str(src))
            before = len(prs.slides._sldIdLst)  # noqa: SLF001
        except Exception:  # noqa: BLE001
            prs = Presentation(str(src))

        summary: dict = {}
        if rebuild_slides is not None:
            # Rebuild: drop every existing slide, rebuild fresh from the list.
            # python-pptx has no public "remove slide"; clear the sldIdLst and
            # drop each slide part's relationship (by rId) from the package so
            # the saved file contains only the new deck.
            for rel in list(prs.part.rels.values()):
                if rel.reltype == RT.SLIDE:
                    prs.part.drop_rel(rel.rId)
            prs.slides._sldIdLst.clear()  # noqa: SLF001
            for spec in rebuild_slides:
                self._add_slide(prs, spec)
            added = len(rebuild_slides)
            summary = {"mode": "rebuild", "removed": before,
                       "added": added, "total": added}
        else:
            # Surgical bin: edits (by original numbering) first, then removals
            # (by original numbering), then appends. Any subset may be combined.
            edited_idx: list[int] = []
            removed_idx: list[int] = []
            for spec in (edit_slides or []):
                idx = int(spec.get("slide", 1))
                edited_idx.append(idx)
                if 1 <= idx <= len(prs.slides._sldIdLst):  # noqa: SLF001
                    self._edit_one_slide(prs.slides[idx - 1], spec)
            for idx in sorted({int(i) for i in (remove_slides or [])}, reverse=True):
                if 1 <= idx <= len(prs.slides._sldIdLst):  # noqa: SLF001
                    self._remove_slide(prs, idx)
                    removed_idx.append(idx)
            for spec in (add_slides or []):
                self._add_slide(prs, spec)
            added = len(add_slides or [])
            mode = "append" if not edited_idx and not removed_idx else "edit_slides"
            summary = {
                "mode": mode,
                "removed": len(removed_idx),
                "added": added,
                "edited": sorted(set(edited_idx)),
                "removed_idx": sorted(removed_idx),
                "total": len(prs.slides._sldIdLst),  # noqa: SLF001
            }

        prs.save(str(dst))
        summary["path"] = str(dst)
        summary["inplace"] = inplace
        return str(dst), summary

    # ------------------------------------------------------------------
    # Reading
    # ------------------------------------------------------------------
    def read_presentation(self, path: str | Path) -> dict:
        src = Path(path)
        if not src.exists():
            raise FileNotFoundError(src)
        prs = Presentation(str(src))
        slides = []
        for slide in prs.slides:
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        t = "".join(run.text for run in p.runs)
                        if t.strip():
                            texts.append(t)
                if shape.has_table:
                    for row in shape.table.rows:
                        texts.append(" | ".join(c.text for c in row.cells))
            slides.append({"index": slide.slide_id, "texts": texts})
        return {"path": str(src), "slide_count": len(slides), "slides": slides}

    def summarize(self, path: str | Path) -> dict:
        data = self.read_presentation(path)
        return {
            "path": data["path"],
            "slide_count": data["slide_count"],
            "first_title": data["slides"][0]["texts"][0] if data["slides"] and data["slides"][0]["texts"] else None,
            "slide_titles": self._titles(data),
        }

    def _titles(self, data: dict) -> list:
        titles = []
        for s in data["slides"]:
            t = s["texts"][0] if s["texts"] else ""
            titles.append(t)
        return titles

    # -- live editing via PowerPoint COM (presentation open in PowerPoint) --
    def is_open_in_ppt(self, path: str | Path) -> bool:
        """Return True if the .pptx is currently open in Microsoft PowerPoint."""
        try:
            return bool(_live.is_open_in_ppt(path))
        except Exception:  # noqa: BLE001
            return False

    def live_edit(self, path: str | Path, add_slides=None, rebuild_slides=None,
                  edit_slides=None, remove_slides=None) -> dict:
        """Append, rebuild, surgically edit, or remove slides in a presentation
        open in PowerPoint, live. Returns
        {'ok','path','mode':'live-ppt','summary':{...}}.
        """
        return _live.live_edit(path, add_slides=add_slides,
                               rebuild_slides=rebuild_slides,
                               edit_slides=edit_slides,
                               remove_slides=remove_slides)

    def live_append(self, path: str | Path, text: str) -> dict:
        """Add a textbox to the last slide, live on screen."""
        return _live.append_text(path, text)
