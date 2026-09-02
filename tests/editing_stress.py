"""
Editing stress suite — the list of edits Varan must be able to perform, tested
through the REAL ToolExecutor (the exact tool surface the agent uses).

Covers every supported document type with complex editing tasks:
  Word (.docx): create, styled content-append, replace/insert/delete,
                delete_range (+end_level), remove_blank_pages, track changes,
                in-place writes, destructive-confirmation gate, get_paragraphs.
  Excel (.xlsx): create, cell writes, formulas, new sheet, charts, edit in place.
  PowerPoint (.pptx): create, read back, append slides (incl. table/chart).
  PDF (.pdf): create fixture, read/summarize, best-effort replace/delete.
  Text (.txt/.md): replace (count and all), insert before/after, append,
                   delete, delete_range — plus not-found errors surface clearly.
  Cross-cutting: create_document-on-open-target remaps to edit_document with
                 structured 'content' (no markdown flattening).

Run:  python tests/editing_stress.py
Exits non-zero if any edit fails, printing a per-spec report.
"""
from __future__ import annotations

import shutil
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

from docx import Document  # noqa: E402

from agent.tools import ToolExecutor  # noqa: E402

WORK = ROOT / "outputs" / "_stress"
shutil.rmtree(WORK, ignore_errors=True)
WORK.mkdir(parents=True, exist_ok=True)

ex = ToolExecutor(WORK, strict=True)

FAILURES: list[str] = []
COUNT = {"pass": 0, "fail": 0}


def spec(name: str):
    def deco(fn):
        def wrapper(*args, **kwargs):
            try:
                fn(*args, **kwargs)
                COUNT["pass"] += 1
                print(f"  [PASS] {name}")
            except Exception as exc:  # noqa: BLE001
                COUNT["fail"] += 1
                FAILURES.append(f"{name}: {type(exc).__name__}: {exc}")
                print(f"  [FAIL] {name} -> {type(exc).__name__}: {exc}")
        return wrapper
    return deco


def call(name: str, args: dict) -> dict:
    """Execute a tool like the agent would; raise on tool-level errors."""
    result = ex.execute(name, args)
    if not isinstance(result, dict):
        raise AssertionError(f"{name} returned non-dict: {result!r}")
    return result


def word_docx(tag: str, body: list | None = None, title: str = "Varan Stress") -> str:
    p = WORK / f"{tag}.docx"
    r = call("create_document", {"path": str(p), "title": title, "body": body or []})
    assert r.get("ok"), r
    return str(Path(r["path"]).resolve())


def docx_texts(p: str | Path) -> list[str]:
    d = Document(str(p))
    return [par.text for par in d.paragraphs]


def docx_paras(p: str | Path):
    d = Document(str(p))
    return [{"text": par.text, "style": par.style.name if par.style is not None else ""} for par in d.paragraphs]


# ======================================================================
# WORD
# ======================================================================
@spec("W1 create styled document (headings/italic/bold/bullets/numbered/table)")
def w1():
    path = word_docx("w1", [
        {"type": "heading", "level": 1, "text": "Chapter 1"},
        {"type": "heading", "level": 3, "text": "Section A"},
        {"type": "paragraph", "text": "Italic intro line.", "italic": True},
        {"type": "paragraph", "text": "A bold claim.", "bold": True},
        {"type": "bullet", "text": "B1"},
        {"type": "bullet", "text": "B2"},
        {"type": "numbered", "text": "N1"},
        {"type": "table", "headers": ["Name", "Score"], "rows": [["Ada", 99], ["Grace", 88]]},
    ])
    paras = docx_paras(path)
    styles = {p["style"] for p in paras}
    assert "Heading 1" in styles and "Heading 3" in styles, styles
    assert "List Bullet" in styles and "List Number" in styles, styles
    assert any("Italic intro line." in p["text"] for p in paras)
    d = Document(path)
    assert len(d.tables) == 1 and d.tables[0].cell(1, 0).text == "Ada"


@spec("W2 replace an exact phrase, keeping surrounding content")
def w2():
    path = word_docx("w2", [
        {"type": "paragraph", "text": "The product ships in June."},
        {"type": "paragraph", "text": "Keep this.", "italic": True},
    ])
    r = call("edit_document", {"path": path, "edits": [
        {"action": "replace", "match": "ships in June", "text": "launches in July"}],
        "inplace": True})
    assert r.get("ok"), r
    texts = docx_texts(path)
    assert any("launches in July" in t for t in texts), texts
    assert any("Keep this." in t for t in texts), "replace dropped unrelated content"


@spec("W3 replace of missing text returns a clear error (no silent ok)")
def w3():
    path = word_docx("w3", [{"type": "paragraph", "text": "Lone paragraph."}])
    r = call("edit_document", {"path": path, "edits": [
        {"action": "replace", "match": "this text does not exist", "text": "x"}],
        "inplace": True})
    assert not r.get("ok"), f"silent success: {r}"
    assert "not found" in str(r.get("error", "")), r


@spec("W4 insert text before and after a matched paragraph")
def w4():
    path = word_docx("w4", [
        {"type": "paragraph", "text": "MIDDLE"},
        {"type": "paragraph", "text": "Tail."},
    ])
    r = call("edit_document", {"path": path, "edits": [
        {"action": "insert_before", "match": "MIDDLE", "text": "HEAD"},
        {"action": "insert_after", "match": "MIDDLE", "text": "AFTER"},
    ], "inplace": True})
    assert r.get("ok"), r
    texts = docx_texts(path)
    hi = texts.index("HEAD")
    mi = texts.index("MIDDLE")
    ai = texts.index("AFTER")
    assert hi < mi < ai, texts


@spec("W5 append plain text at the end")
def w5():
    path = word_docx("w5", [{"type": "paragraph", "text": "Start"}])
    r = call("edit_document", {"path": path, "append": "END LINE", "inplace": True})
    assert r.get("ok"), r
    assert docx_texts(path)[-1].strip() == "END LINE"


@spec("W6 styled content-append renders REAL Word styles (not '##' markdown)")
def w6():
    path = word_docx("w6", [{"type": "paragraph", "text": "Start"}])
    r = call("edit_document", {"path": path, "content": [
        {"type": "heading", "level": 1, "text": "IRON MAN"},
        {"type": "heading", "level": 2, "text": "Origins"},
        {"type": "paragraph", "text": "Cave scene.", "courier": True, "italic": True},
        {"type": "bullet", "text": "First suit"},
        {"type": "divider"},
        {"type": "paragraph", "text": "Bold lead.", "bold": True},
    ], "inplace": True})
    assert r.get("ok"), r
    paras = docx_paras(path)
    full = "\n".join(p["text"] for p in paras)
    assert "## " not in full and not any(p["text"].startswith("- ") for p in paras), "markdown leaked into docx"
    tail = [p for p in paras if p["text"].strip() == "IRON MAN"]
    assert tail and tail[0]["style"].startswith("Heading"), tail
    h2 = [p for p in paras if p["style"] == "Heading 2" and p["text"].strip() == "Origins"]
    assert h2, paras
    lists = [p for p in paras if p["style"] in ("List Bullet", "List Number")]
    assert lists, paras
    assert any("* * *" in p["text"] for p in paras), "scene-break divider missing"


@spec("W7 damage control: unstyled 'append' does NOT leak markdown when remap is bypassed")
def w7():
    path = word_docx("w7", [{"type": "paragraph", "text": "Start"}])
    r = call("edit_document", {"path": path, "content": [
        {"type": "heading", "level": 1, "text": "# RAW HEAD"}], "inplace": True})
    assert r.get("ok"), r
    texts = docx_texts(path)
    assert "## " not in "\n".join(texts), "markdown leaked"


@spec("W8 destructive delete REFUSED without confirm (need_confirmation)")
def w8():
    path = word_docx("w8", [
        {"type": "heading", "level": 1, "text": "ARC 1"},
        {"type": "paragraph", "text": "KEEP ME"},
    ])
    r = call("edit_document", {"path": path, "edits": [
        {"action": "delete", "match": "ARC 1"}], "inplace": True})
    assert r.get("need_confirmation") is True, f"delete allowed without confirm: {r}"
    assert "KEEP ME" in "\n".join(docx_texts(path)), "doc was modified despite refusal"


@spec("W9 destructive delete performed AFTER confirm (removes the matched paragraph)")
def w9():
    path = word_docx("w9", [
        {"type": "heading", "level": 1, "text": "ARC 2"},
        {"type": "paragraph", "text": "KEEP"},
    ])
    r = call("edit_document", {"path": path, "edits": [
        {"action": "delete", "match": "ARC 2"}], "confirm": True, "inplace": True})
    assert r.get("ok"), r
    texts = docx_texts(path)
    assert "ARC 2" not in "\n".join(texts), texts
    assert "KEEP" in "\n".join(texts)


@spec("W10 delete_range between anchors (confirmed)")
def w10():
    path = word_docx("w10", [
        {"type": "heading", "level": 1, "text": "Section Start"},
        {"type": "paragraph", "text": "Middle filler"},
        {"type": "paragraph", "text": "Section End"},  # last occurrence
        {"type": "paragraph", "text": "SURVIVES"},
    ])
    r = call("edit_document", {"path": path, "edits": [
        {"action": "delete_range", "match": "Section Start", "end_match": "Section End"}],
        "confirm": True, "inplace": True})
    assert r.get("ok"), r
    texts = docx_texts(path)
    assert "SURVIVES" in "\n".join(texts), texts
    assert "Section Start" not in "\n".join(texts) and "Middle filler" not in "\n".join(texts)


@spec("W11 delete_range with end_level stops before the next Heading 1")
def w11():
    path = word_docx("w11", [
        {"type": "heading", "level": 1, "text": "Delete Me"},
        {"type": "paragraph", "text": "Part of doomed block"},
        {"type": "heading", "level": 1, "text": "Next Chapter"},
        {"type": "paragraph", "text": "SURVIVES"},
    ])
    r = call("edit_document", {"path": path, "edits": [
        {"action": "delete_range", "match": "Delete Me", "end_level": 1}],
        "confirm": True, "inplace": True})
    assert r.get("ok"), r
    texts = docx_texts(path)
    assert "Delete Me" not in "\n".join(texts) and "Part of doomed block" not in "\n".join(texts), texts
    assert "Next Chapter" in "\n".join(texts) and "SURVIVES" in "\n".join(texts), texts


@spec("W12 remove_blank_pages removes only empty paragraphs and is idempotent")
def w12():
    path = word_docx("w12", [  # python-docx leaves one trailing empty para; add real empties
        {"type": "heading", "level": 1, "text": "Keep"},
        {"type": "paragraph", "text": "Body text"},
    ])
    d = Document(path)
    for _ in range(5):
        d.add_heading("", level=1)
    d.save(path)
    r = call("edit_document", {"path": path, "remove_blank_pages": True, "inplace": True})
    assert r.get("ok") and r.get("removed", 0) >= 5, r
    remaining = [p for p in docx_paras(path) if not p["text"].strip()]
    assert remaining == [], f"empties remain: {remaining}"
    r2 = call("edit_document", {"path": path, "remove_blank_pages": True, "inplace": True})
    assert r2.get("removed") == 0, f"second run not a no-op: {r2}"


@spec("W13 get_paragraphs reports empty-paragraph count (blank-page visibility)")
def w13():
    path = word_docx("w13", [{"type": "paragraph", "text": "Real line"}])
    d = Document(path)
    for _ in range(3):
        d.add_heading("", level=1)
    d.save(path)
    r = call("get_paragraphs", {"path": path})
    assert r.get("ok"), r
    joined = "\n".join(r.get("paragraphs", []))
    assert "-- 3 empty paragraph(s)" in joined or "3 empty paragraph(s)" in joined, r


@spec("W14 tracked changes (redlines) write valid OOXML (w:del + w:ins)")
def w14():
    path = word_docx("w14", [{"type": "paragraph", "text": "Old sentence here."}])
    r = call("edit_document", {"path": path, "edits": [
        {"action": "replace", "match": "Old sentence", "text": "NEW sentence"}],
        "track_changes": True, "author": "Ada", "inplace": True})
    assert r.get("ok"), r
    import zipfile
    with zipfile.ZipFile(path) as zf:
        xml = zf.read("word/document.xml").decode("utf-8")
    assert "<w:del " in xml and ">Old sentence" in xml, "deletion mark missing"
    assert "<w:ins " in xml and "NEW sentence" in xml, "insertion mark missing"


@spec("W15 in-place edit on the target leaves NO '_edited' duplicate")
def w15():
    path = word_docx("w15", [{"type": "paragraph", "text": "Change me"}])
    before = set(WORK.glob("w15*"))
    r = call("edit_document", {"path": path, "edits": [
        {"action": "replace", "match": "Change me", "text": "Changed"}],
        "inplace": True})
    assert r.get("ok"), r
    after = set(WORK.glob("w15*"))
    assert after == before, f"duplicate produced: {after - before}"


@spec("W16 read_file + summarize_file work on the edited doc")
def w16():
    path = word_docx("w16", [{"type": "heading", "level": 1, "text": "Read Me"},
                             {"type": "paragraph", "text": "some rich body text"}])
    r = call("read_file", {"path": path})
    assert r.get("ok"), r
    assert "Read Me" in str(r.get("data", {})), r
    s = call("summarize_file", {"path": path})
    assert s.get("ok"), s
    summ = str(s.get("summary", ""))
    assert "Read Me" in summ


@spec("W17 replace ALL occurrences across the document (count=-1, formatting-preserving)")
def w17():
    src = word_docx("w17", [
        {"type": "paragraph", "text": "the word dog in the first sentence"},
        {"type": "paragraph", "text": "another dog here"},
        {"type": "heading", "level": 2, "text": "dogs are great"},
    ])
    r = call("edit_document", {"path": src, "edits": [
        {"action": "replace", "match": "dog", "text": "cat", "count": -1}], "inplace": True})
    assert r.get("ok"), r
    paras = docx_paras(src)
    joined = " | ".join(p["text"] for p in paras)
    assert joined.count("cat") == 3, paras          # body(2) + heading(1)
    assert "dog" not in joined, paras
    heading = [p for p in paras if p["style"].startswith("Heading")]
    assert heading and "cats are great" in heading[0]["text"], paras  # style preserved


@spec("W18 complex table edit: set specific cell text inside an existing table")
def w18():
    src = word_docx("w18", [
        {"type": "paragraph", "text": "intro"},
        {"type": "table", "headers": ["Name", "Role"], "rows": [
            ["Alice", "Dev"], ["Bob", "QA"]]},
    ])
    r = call("edit_document", {"path": src, "table_edits": [
        {"table": 1, "cell": {"ref": "B2"}, "text": "Lead Dev"}], "inplace": True})
    assert r.get("ok"), r
    data = call("read_file", {"path": src})
    blob = str(data.get("data"))
    assert "Lead Dev" in blob, data
    assert "Alice" in blob and "Bob" in blob and "Dev" in blob, data  # other cells untouched
    assert "QA" in blob and "Lead Dev" in blob  # QA stays in Bob's cell (B3), only B2 changed


# ======================================================================
# EXCEL
# ======================================================================
@spec("X1 create workbook with headers + rows")
def x1():
    p = WORK / "x1.xlsx"
    r = call("create_workbook", {"path": str(p), "sheet": "Sales", "data": [
        {"cells": ["Month", "Revenue"]},
        {"cells": ["Jan", 100]},
        {"cells": ["Feb", 150]},
    ]})
    assert r.get("ok"), r
    rb = call("read_file", {"path": str(p)})
    assert rb.get("ok") and "Sales" in str(rb.get("data")), rb


@spec("X2 edit_workbook writes cells, formulas, new sheet, and a chart")
def x2():
    p = WORK / "x2.xlsx"
    call("create_workbook", {"path": str(p), "sheet": "Data", "data": [
        {"cells": ["Q1", "Q2"]},
        {"cells": [10, 20]},
        {"cells": [30, 40]},
    ]})
    r = call("edit_workbook", {"path": str(p), "writes": [
        {"cell": "C1", "value": "Total"},
        {"cell": "C2", "value": "=SUM(A2:B2)"},
    ], "formulas": [{"cell": "C3", "formula": "=SUM(A3:B3)"}],
        "new_sheet": "Totals", "add_chart": {
            "type": "bar", "title": "Revenue", "categories": "A2:A3", "data": "B2:B3"},
        "inplace": True})
    assert r.get("ok"), r
    from openpyxl import load_workbook
    wb = load_workbook(p, data_only=False)
    ws = wb["Data"]
    assert ws["C1"].value == "Total" and ws["C2"].value == "=SUM(A2:B2)", (ws["C1"].value, ws["C2"].value)
    assert ws["C3"].value == "=SUM(A3:B3)", ws["C3"].value
    assert "Totals" in wb.sheetnames, wb.sheetnames
    assert len(ws._charts) >= 1, "chart missing"


@spec("X3 in-place workbook edit leaves no duplicate")
def x3():
    p = WORK / "x3.xlsx"
    call("create_workbook", {"path": str(p), "sheet": "S", "data": [{"cells": ["A"]}, {"cells": [1]}]})
    before = set(WORK.glob("x3*"))
    r = call("edit_workbook", {"path": str(p), "writes": [
        {"cell": "B1", "value": 42}], "inplace": True})
    assert r.get("ok"), r
    assert set(WORK.glob("x3*")) == before, "duplicate produced"


@spec("X4 complex edits: insert/delete rows+columns, clear ranges, styles, delete sheet")
def x4():
    from openpyxl import load_workbook
    p = WORK / "x4.xlsx"
    call("create_workbook", {"path": str(p), "sheet": "Alpha", "data": [
        {"cells": ["Name", "Value", "Extra"]},
        {"cells": ["a", 1, "keep"]},
        {"cells": ["b", 2, "keep"]},
        {"cells": ["c", 3, "keep"]},
    ]})
    r = call("edit_workbook", {
        "path": str(p),
        "new_sheet": "ToRemove",
        "columns": [{"action": "delete", "at": "C"}],       # drop the 'Extra' column
        "rows": [{"action": "insert", "at": 2},             # blank row above 'a'
                 {"action": "delete", "at": 4}],            # delete 'b' row
        "clear": ["A2:A2"],                                 # clear 'Name' col row 2 (blank row)
        "styles": [{"cell": "A1", "bold": True, "fill": "FFD966"}],
        "delete_sheet": "ToRemove",
        "inplace": True,
    })
    assert r.get("ok"), r
    wb = load_workbook(p)
    assert "ToRemove" not in wb.sheetnames, wb.sheetnames   # sheet deleted
    ws = wb["Alpha"]
    assert ws.max_column == 2, ws.max_column                # column C removed
    # row 1 = headers, row 2 = inserted blank, row 3 = a, row 4 = c (b deleted)
    labels = [ws.cell(row=i, column=1).value for i in range(1, ws.max_row + 1)]
    assert labels == ["Name", None, "a", "c"], labels
    assert ws["A1"].font.bold, "header not bold"
    assert ws["A1"].fill.start_color.rgb in ("00FFD966", "FFD966"), ws["A1"].fill.start_color.rgb


# ======================================================================
# POWERPOINT
# ======================================================================
@spec("P1 create presentation with title/bullets/blank+table layouts")
def p1():
    p = WORK / "p1.pptx"
    r = call("create_presentation", {"path": str(p), "slides": [
        {"layout": "title", "title": "Deck", "subtitle": "sub"},
        {"layout": "bullets", "title": "Agenda", "bullets": ["A", "B"]},
        {"layout": "title", "title": "Table Slide", "table": {
            "headers": ["Name", "Role"], "rows": [["Ada", "PM"], ["Grace", "Eng"]]}},
    ]})
    assert r.get("ok"), r
    rb = call("read_file", {"path": str(p)})
    assert rb.get("ok") and "Deck" in str(rb.get("data")), rb


@spec("P2 append slides to existing deck (count grows, content present)")
def p2():
    p = WORK / "p2.pptx"
    call("create_presentation", {"path": str(p), "slides": [
        {"layout": "title", "title": "First"}]})
    r = call("edit_presentation", {"path": str(p), "add_slides": [
        {"layout": "bullets", "title": "Second Deck", "bullets": ["x", "y", "z"]},
        {"layout": "title", "title": "Third Slide"},
    ], "inplace": True})
    assert r.get("ok"), r
    assert r.get("summary", {}).get("mode") == "append", r
    assert r["summary"]["removed"] == 0 and r["summary"]["added"] == 2, r
    rb = call("read_file", {"path": str(p)})
    data = str(rb.get("data", ""))
    assert "First" in data and "Second Deck" in data and "Third Slide" in data, rb


@spec("P3 REBUILD replaces the whole deck (retheme) — mode=rebuild, old slides gone")
def p3():
    from pptx import Presentation
    p = WORK / "p3.pptx"
    call("create_presentation", {"path": str(p), "slides": [
        {"layout": "title", "title": "OLD TITLE", "subtitle": "X"},
        {"layout": "bullets", "title": "OLD SLIDE 2", "bullets": ["a", "b"]},
        {"layout": "title", "title": "OLD SLIDE 3", "subtitle": "Y"},
    ]})
    assert len(Presentation(str(p)).slides) == 3
    r = call("edit_presentation", {"path": str(p), "rebuild_slides": [
        {"layout": "title", "title": "NEW TITLE", "subtitle": "Fresh"},
        {"layout": "bullets", "title": "Market Demand", "bullets": ["x", "y"]},
    ], "inplace": True})
    assert r.get("ok"), r
    summ = r.get("summary", {})
    assert summ.get("mode") == "rebuild", r
    assert summ.get("removed") == 3 and summ.get("added") == 2 and summ.get("total") == 2, r
    prs = Presentation(str(p))
    assert len(prs.slides) == 2, f"expect 2 slides after rebuild, got {len(prs.slides)}"
    body = " ".join(sh.text_frame.text for s in prs.slides for sh in s.shapes if sh.has_text_frame)
    assert "OLD" not in body and "NEW TITLE" in body and "Market Demand" in body, body


@spec("P4 SURGICAL edit_slides changes only specific slides (title + replace + bullets) without touching the rest")
def p4():
    from pptx import Presentation
    p = WORK / "p4.pptx"
    call("create_presentation", {"path": str(p), "slides": [
        {"layout": "title", "title": "Deck Title", "subtitle": "Sub"},
        {"layout": "bullets", "title": "Overview", "bullets": ["Point one", "Point two"]},
        {"layout": "title", "title": "Closing", "subtitle": "Thanks"},
    ]})
    assert len(Presentation(str(p)).slides) == 3
    r = call("edit_presentation", {"path": str(p), "edit_slides": [
        {"slide": 1, "title": "Deck Title v2"},
        {"slide": 2, "replace": {"Point one": "Point one (updated)"}},
    ], "inplace": True})
    assert r.get("ok"), r
    summ = r.get("summary", {})
    assert summ.get("mode") == "edit_slides", r
    assert summ.get("edited") == [1, 2], r
    assert summ.get("total") == 3, r
    prs = Presentation(str(p))
    assert len(prs.slides) == 3, f"expect unchanged slide count, got {len(prs.slides)}"
    texts = [sh.text_frame.text for s in prs.slides for sh in s.shapes if sh.has_text_frame]
    joined = " | ".join(texts)
    assert "Deck Title v2" in joined, joined
    assert "Point one (updated)" in joined, joined
    assert "Point two" in joined, joined   # untouched bullet preserved
    assert "Closing" in joined, joined      # slide 3 untouched
    assert "Deck Title" not in joined.replace("Deck Title v2", ""), joined


@spec("P5 create_presentation on an open target routes to REBUILD (not append) so a new deck never duplicates the old one")
def p5():
    from agent.loop import Agent
    ag = Agent.__new__(Agent)
    ag.target_file = str(WORK / "p5_target.pptx")
    remapped = ag._remap_create_to_edit("create_presentation", {
        "path": "elsewhere.pptx",
        "slides": [
            {"layout": "title", "title": "Fresh Deck", "subtitle": ""},
            {"layout": "bullets", "title": "Topic", "bullets": ["a"]},
        ],
    })
    assert remapped is not None, "no remap happened"
    name, args = remapped
    assert name == "edit_presentation", name
    assert args.get("rebuild_slides") is not None, args
    assert args.get("add_slides") is None, args
    assert args.get("path") == ag.target_file, args


@spec("P6 remove_slides deletes SPECIFIC slides surgically (no rebuild, others untouched)")
def p6():
    from pptx import Presentation
    p = WORK / "p6.pptx"
    call("create_presentation", {"path": str(p), "slides": [
        {"layout": "title", "title": "One", "subtitle": ""},
        {"layout": "bullets", "title": "Two", "bullets": ["b"]},
        {"layout": "bullets", "title": "Three", "bullets": ["c"]},
        {"layout": "title", "title": "Four", "subtitle": ""},
    ]})
    assert len(Presentation(str(p)).slides) == 4
    r = call("edit_presentation", {"path": str(p), "remove_slides": [2, 4], "inplace": True})
    assert r.get("ok"), r
    summ = r.get("summary", {})
    assert summ.get("removed") == 2, r
    assert summ.get("removed_idx") == [2, 4], r
    assert summ.get("total") == 2, r
    prs = Presentation(str(p))
    assert len(prs.slides) == 2, len(prs.slides)
    titles = []
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_text_frame and (sh.text_frame.text or "").strip():
                titles.append(sh.text_frame.text)
    joined = " | ".join(titles)
    assert "One" in joined and "Three" in joined, joined
    assert "Two" not in joined and "Four" not in joined, joined


@spec("P7 SURGICAL add shapes on one slide: add_textbox + add_table + replace in table cells")
def p7():
    from pptx import Presentation
    p = WORK / "p7.pptx"
    call("create_presentation", {"path": str(p), "slides": [
        {"layout": "title", "title": "Top", "subtitle": ""},
        {"layout": "bullets", "title": "Mid", "bullets": ["hello"]},
    ]})
    r = call("edit_presentation", {"path": str(p), "edit_slides": [
        {"slide": 2, "add_textbox": {"text": "ANNOTATION"},
         "add_table": {"headers": ["K", "V"], "rows": [["a", "b"], ["c", "d"]]}},
        {"slide": 2, "replace": {"hello": "world"}},
    ], "inplace": True})
    assert r.get("ok"), r
    prs = Presentation(str(p))
    assert len(prs.slides) == 2, len(prs.slides)  # no slides appended
    slide2 = prs.slides[1]
    texts = " | ".join(sh.text_frame.text for sh in slide2.shapes if sh.has_text_frame)
    assert "world" in texts, texts
    assert "ANNOTATION" in texts, texts
    tables = [sh for sh in slide2.shapes if sh.has_table]
    assert len(tables) == 1, "table not added"
    grid = [[c.text for c in row.cells] for row in tables[0].table.rows]
    assert grid[0] == ["K", "V"] and grid[1] == ["a", "b"], grid


@spec("P8 placeholder filler content is REFUSED (no silent '- item' success)")
def p8():
    p = WORK / "p8.pptx"
    r = call("create_presentation", {"path": str(p), "slides": [
        {"layout": "title", "title": "Bad", "subtitle": ""},
        {"layout": "bullets", "title": "Filler", "bullets": ["item"]},
    ]})
    assert "error" in r and "Refusing" in r["error"], r
    assert not p.exists(), "refused create must not write a file"
    r2 = call("create_presentation", {"path": str(p), "slides": [
        {"layout": "bullets", "title": "BadTable", "bullets": ["real"],
         "table": {"headers": ["item"], "rows": [["i", "t", "e", "m"]]}},
    ]})
    assert "error" in r2 and "Refusing" in r2["error"], r2
    assert not p.exists(), "refused create must not write a file"
    # surgical filler is refused the same way
    good = WORK / "p8_good.pptx"
    call("create_presentation", {"path": str(good), "slides": [
        {"layout": "bullets", "title": "Good", "bullets": ["one", "two"]},
    ]})
    r3 = call("edit_presentation", {"path": str(good), "edit_slides": [
        {"slide": 1, "set_bullets": ["item"]}], "inplace": True})
    assert "error" in r3 and "Refusing" in r3["error"], r3


# ======================================================================
# PDF
# ======================================================================
@spec("D1 read/summarize a generated PDF")
def d1():
    from reportlab.pdfgen import canvas as _canvas
    p = WORK / "d1.pdf"
    c = _canvas.Canvas(str(p))
    c.setFont("Helvetica", 12)
    c.drawString(72, 780, "SECTION ALPHA")
    c.drawString(72, 760, "First body line here.")
    c.save()
    s = call("summarize_file", {"path": str(p)})
    assert s.get("ok") and "pdf" in str(s.get("summary", "").get("kind", "")), s
    r = call("read_file", {"path": str(p)})
    assert r.get("ok") and "SECTION ALPHA" in str(r.get("data")), r


@spec("D2 best-effort replace text in a simple PDF")
def d2():
    from reportlab.pdfgen import canvas as _canvas
    p = WORK / "d2.pdf"
    c = _canvas.Canvas(str(p))
    c.setFont("Helvetica", 12)
    c.drawString(72, 780, "SECTION ALPHA")
    c.drawString(72, 760, "First body line here.")
    c.save()
    r = call("edit_text", {"path": str(p), "edits": [
        {"action": "replace", "match": "First body line here.", "text": "EDITED LINE"}],
        "inplace": True})
    assert r.get("ok"), r
    from pypdf import PdfReader
    txt = PdfReader(p).pages[0].extract_text() or ""
    assert "EDITED LINE" in txt, txt


# ======================================================================
# TEXT (.md / .txt)
# ======================================================================
@spec("T1 replace first and Nth occurrence (count)")
def t1():
    p = WORK / "t1.md"
    p.write_text("one two\none two\none two\n", encoding="utf-8")
    r = call("edit_text", {"path": str(p), "edits": [
        {"action": "replace", "match": "one", "text": "uno"}], "inplace": True})
    assert r.get("ok"), r
    assert p.read_text(encoding="utf-8").count("uno") == 1
    r = call("edit_text", {"path": str(p), "edits": [
        {"action": "replace", "match": "one", "text": "uno", "count": 2}], "inplace": True})
    assert r.get("ok"), r
    assert p.read_text(encoding="utf-8").count("uno") == 3


@spec("T2 replace ALL occurrences (count=-1)")
def t2():
    p = WORK / "t2.md"
    p.write_text("cat dog\ncat dog\ncat bird\n", encoding="utf-8")
    r = call("edit_text", {"path": str(p), "edits": [
        {"action": "replace", "match": "cat", "text": "CAT", "count": -1}], "inplace": True})
    assert r.get("ok"), r
    assert "CAT dog\nCAT dog\nCAT bird\n" == p.read_text(encoding="utf-8")


@spec("T3 insert before and after a matched line")
def t3():
    p = WORK / "t3.md"
    p.write_text("MID\n", encoding="utf-8")
    r = call("edit_text", {"path": str(p), "edits": [
        {"action": "insert_before", "match": "MID", "text": "==BEFORE==\n"},
        {"action": "insert_after", "match": "MID", "text": "\n==AFTER=="},
    ], "inplace": True})
    assert r.get("ok"), r
    assert p.read_text(encoding="utf-8") == "==BEFORE==\nMID\n==AFTER==\n"


@spec("T4 append text at the end")
def t4():
    p = WORK / "t4.md"
    p.write_text("head\n", encoding="utf-8")
    r = call("edit_text", {"path": str(p), "append": "tail line", "inplace": True})
    assert r.get("ok"), r
    assert p.read_text(encoding="utf-8").strip().endswith("tail line")


@spec("T5 delete a line")
def t5():
    p = WORK / "t5.md"
    p.write_text("keep1\ndoomed line\nkeep2\n", encoding="utf-8")
    r = call("edit_text", {"path": str(p), "edits": [
        {"action": "delete", "match": "doomed line"}], "inplace": True})
    assert r.get("ok"), r
    assert "doomed line" not in p.read_text(encoding="utf-8")
    assert "keep1" in p.read_text(encoding="utf-8") and "keep2" in p.read_text(encoding="utf-8")


@spec("T6 delete_range removes a block between anchors")
def t6():
    p = WORK / "t6.md"
    p.write_text("# START\ncontent\n## END\nsurvivor\n", encoding="utf-8")
    r = call("edit_text", {"path": str(p), "edits": [
        {"action": "delete_range", "match": "# START", "end_match": "## END"}],
        "inplace": True})
    assert r.get("ok"), r
    assert p.read_text(encoding="utf-8") == "survivor\n"


@spec("T7 missing text returns a clear 'not found' error")
def t7():
    p = WORK / "t7.md"
    p.write_text("hello world\n", encoding="utf-8")
    r = call("edit_text", {"path": str(p), "edits": [
        {"action": "replace", "match": "absent", "text": "x"}], "inplace": True})
    assert not r.get("ok"), r
    assert "not found" in str(r.get("error", "")).lower(), r


# ======================================================================
# CROSS-CUTTING: create_document on an open target -> edit with content
# ======================================================================
@spec("C1 create_document on the open target REMAPS to edit_document with structured content")
def c1():
    from agent.loop import Agent
    target = WORK / "c1.docx"
    call("create_document", {"path": str(target), "body": [
        {"type": "paragraph", "text": "Existing"}]})
    ag = Agent(None, {})  # type: ignore[arg-type]
    ag.target_file = str(target)
    remapped = ag._remap_create_to_edit("create_document", {
        "title": "NEW SECTION",
        "body": [
            {"type": "heading", "level": 2, "text": "H2 Title"},
            {"type": "paragraph", "text": "Styled body.", "italic": True},
        ],
    })
    assert remapped is not None, "no remap happened"
    name, args = remapped
    assert name == "edit_document", name
    assert "append" not in args and "content" in args, f"stale markdown flatten: {args}"
    assert args["content"][0]["type"] == "title", args


@spec("R1 read_file returns clean anydoc Markdown when the optional package is installed")
def r1():
    from office import anydoc_reader
    p = WORK / "r1.docx"
    call("create_document", {"path": str(p), "body": [
        {"type": "heading", "level": 1, "text": "AnyDoc Heading"},
        {"type": "paragraph", "text": "Body line for markdown."},
        {"type": "table", "headers": ["Col A", "Col B"], "rows": [["1", "2"]]},
    ]})
    r = call("read_file", {"path": str(p)})
    assert r.get("ok") and "AnyDoc Heading" in str(r.get("data")), r
    if anydoc_reader.available():
        md = r["data"].get("markdown")
        assert md, "anydoc installed but no markdown attached"
        assert "AnyDoc Heading" in md and "Col A" in md, md
    else:
        # graceful fallback: structured read still works without anydoc
        assert "paragraphs" in r["data"], r


def main() -> int:
    print("Varan editing stress suite — exercising the real ToolExecutor surface\n")
    _ = (w1, w2, w3, w4, w5, w6, w7, w8, w9, w10, w11, w12, w13, w14, w15, w16,
         w17, w18, x1, x2, x3, x4, p1, p2, p3, p4, p5, p6, p7, p8, d1, d2,
         t1, t2, t3, t4, t5, t6, t7, c1, r1)
    for fn in _:
        fn()
    print(f"\n{COUNT['pass']} passed, {COUNT['fail']} failed")
    if FAILURES:
        print("\nFAILED EDITS:")
        for f in FAILURES:
            print(f"  - {f}")
        return 1
    print("\n[PASS] Every edit Varan must be able to perform works.")
    return 0


if __name__ == "__main__":
    sys.exit(main())